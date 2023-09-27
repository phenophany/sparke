import os
from typing import Optional

import docx2txt
import mammoth
import pypandoc
import re
import fitz
from pptx import Presentation


# 将各类文件类型转换为txt文本
def wordtohtml(sourcepath:str, destpath:str):

    for root, directories, files in os.walk(sourcepath, topdown=False):
        for file in files:
            if file.lower().endswith(".docx"):
                name = (os.path.join(root, file))
                # text = docx2txt.process(name)
                # metadata = {"file_name": file.name}
                # pypandoc.convert_file(name, 'markdown', outputfile=name[:-5] + ".txt")

                result = mammoth.convert_to_html(name)
                html = result.value  # The generated HTML
                messages = result.messages  # Any messages, such as warnings during conversion
                # 将html保存到destpath
                with open(destpath + "\\" + file[:-5] + ".html", "w", encoding="utf-8") as f:
                    f.write(html)

def wordtomarkdown(sourcepath:str, destpath:str):

    for root, directories, files in os.walk(sourcepath, topdown=False):
        for file in files:
            if file.lower().endswith(".docx"):
                name = (os.path.join(root, file))
                # text = docx2txt.process(name)
                # metadata = {"file_name": file.name}
                # pypandoc.convert_file(name, 'markdown', outputfile=name[:-5] + ".txt")

                result = mammoth.convert_to_markdown(name)
                md = result.value  # The generated HTML
                messages = result.messages  # Any messages, such as warnings during conversion
                # 将html保存到destpath
                with open(destpath + "\\" + file[:-5] + ".md", "w", encoding="utf-8") as f:
                    f.write(md)

def pdftohtml(pdf_folder:str, html_folder:str):
    # 输入PDF文件夹路径
    # pdf_folder = 'path/to/your/pdf_folder'  # 替换为你的PDF文件夹路径
    #
    # # 输出HTML文件夹路径
    # html_folder = 'path/to/your/html_folder'  # 替换为你的HTML文件夹路径

    # 创建输出HTML文件夹
    os.makedirs(html_folder, exist_ok=True)

    # 遍历PDF文件夹中的所有PDF文件
    for root, _, files in os.walk(pdf_folder):
        for pdf_file in files:
            if pdf_file.lower().endswith('.pdf'):
                pdf_path = os.path.join(root, pdf_file)

                # 打开PDF文件
                pdf_document = fitz.open(pdf_path)

                # 创建一个空的HTML字符串，用于存储所有页面的内容
                html_content = ''

                # 遍历PDF的所有页面
                for page_num in range(pdf_document.page_count):
                    page = pdf_document[page_num]
                    page_html = page.get_text("html")

                    # 将每个页面的HTML内容添加到总的HTML字符串中
                    html_content += page_html

                # 关闭PDF文件
                pdf_document.close()

                html_file = os.path.join(html_folder, pdf_file[:-4] + '.html')
                # 将所有页面的HTML内容保存到单个HTML文件
                with open(html_file, "w", encoding="utf-8") as output_html_file:
                    output_html_file.write(html_content)

                print(f"PDF转换为HTML完成，保存在 {html_file} 中。")
    print("PDF转换为HTML完成。")

def pdftomarkdown(pdf_folder:str, html_folder:str):
    # 输入PDF文件夹路径
    # pdf_folder = 'path/to/your/pdf_folder'  # 替换为你的PDF文件夹路径
    #
    # # 输出HTML文件夹路径
    # html_folder = 'path/to/your/html_folder'  # 替换为你的HTML文件夹路径

    # 创建输出HTML文件夹
    os.makedirs(html_folder, exist_ok=True)

    # 遍历PDF文件夹中的所有PDF文件
    for root, _, files in os.walk(pdf_folder):
        for pdf_file in files:
            if pdf_file.lower().endswith('.pdf'):
                pdf_path = os.path.join(root, pdf_file)

                # 打开PDF文件
                pdf_document = fitz.open(pdf_path)

                # 创建一个空的HTML字符串，用于存储所有页面的内容
                html_content = ''

                # 遍历PDF的所有页面
                for page_num in range(pdf_document.page_count):
                    page = pdf_document[page_num]
                    page_html = page.get_text()

                    # 将每个页面的HTML内容添加到总的HTML字符串中
                    html_content += page_html

                # 关闭PDF文件
                pdf_document.close()

                html_file = os.path.join(html_folder, pdf_file[:-4] + '.md')
                # 将所有页面的HTML内容保存到单个HTML文件
                with open(html_file, "w", encoding="utf-8") as output_html_file:
                    output_html_file.write(html_content)

                print(f"PDF转换为Markdown完成，保存在 {html_file} 中。")

def ppttomarkdown(sourcepath:str,destpath:str):
    for root, directories, files in os.walk(sourcepath, topdown=False):
        for file in files:
            if file.lower().endswith(".pptx"):
                ppt_path = os.path.join(root, file)
                # 打开PPT文件
                presentation = Presentation(ppt_path)
                markdown_file = destpath + "\\" + file[:-5] + ".md"

                # 创建Markdown文件并写入内容
                with open(markdown_file, "w", encoding="utf-8") as md_file:
                    # 遍历每个幻灯片
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text = shape.text
                                md_file.write(text)
                                md_file.write("\n")

ppttomarkdown("D:\\test\\sparke\\ppt","D:\\test\\sparke\\ppt")



